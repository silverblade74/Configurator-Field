# export_utils.py
"""
Tiny helpers that turn the in-memory `bundle` dict from app.py
into Excel, PDF, and PowerPoint byte-streams for download.

• to_excel_bytes(bundle)  → bytes
• to_pdf_bytes(bundle)    → bytes
• to_pptx_bytes(bundle)   → bytes
"""

import io
from typing import Dict

import pandas as pd
from reportlab.lib.pagesizes import LETTER
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt


# ────────────────────────────────────────────────────────────────────────────
# Excel
# ────────────────────────────────────────────────────────────────────────────
def to_excel_bytes(bundle: Dict) -> bytes:
    """Return an .xlsx file in memory."""
    output = io.BytesIO()

    # Flatten dicts into DataFrames
    meta_df = pd.DataFrame(bundle["Metadata"], index=[0])
    capture_df = pd.json_normalize(bundle["Capture"]).T.reset_index()
    capture_df.columns = ["Field", "Value"]
    rec_df = pd.DataFrame(
        list(bundle["Recommendations"].items()), columns=["Domain", "Recommendation"]
    )

    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        meta_df.to_excel(writer, sheet_name="Metadata", index=False)
        capture_df.to_excel(writer, sheet_name="Capture", index=False)
        rec_df.to_excel(writer, sheet_name="Recommendations", index=False)

    return output.getvalue()


# ────────────────────────────────────────────────────────────────────────────
# PDF  (simple 1-page text dump)
# ────────────────────────────────────────────────────────────────────────────
def to_pdf_bytes(bundle: Dict) -> bytes:
    output = io.BytesIO()
    c = canvas.Canvas(output, pagesize=LETTER)
    width, height = LETTER
    x, y = 40, height - 40

    def draw_line(text: str, offset: int = 14):
        nonlocal y
        c.drawString(x, y, text)
        y -= offset

    c.setFont("Helvetica-Bold", 12)
    draw_line("SHI Field Configurator – Assessment Summary", 18)

    c.setFont("Helvetica", 10)
    for k, v in bundle["Metadata"].items():
        draw_line(f"{k}: {v}")

    draw_line("")
    draw_line("Capture", 16)
    for k, v in bundle["Capture"].items():
        draw_line(f"• {k}: {v}")

    draw_line("")
    draw_line("Recommendations", 16)
    for k, v in bundle["Recommendations"].items():
        draw_line(f"• {k}: {v}")

    c.showPage()
    c.save()
    return output.getvalue()


# ────────────────────────────────────────────────────────────────────────────
# PowerPoint  (1 slide)
# ────────────────────────────────────────────────────────────────────────────
def to_pptx_bytes(bundle: Dict) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = "SHI Field Configurator – Summary"
    title_tf.paragraphs[0].font.size = Pt(24)

    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(6))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.margin_bottom = 0

    def add_line(text: str, bold=False):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.bold = bold

    add_line("Metadata", bold=True)
    for k, v in bundle["Metadata"].items():
        add_line(f"• {k}: {v}")

    add_line("")  # spacer
    add_line("Capture", bold=True)
    for k, v in bundle["Capture"].items():
        add_line(f"• {k}: {v}")

    add_line("")  # spacer
    add_line("Recommendations", bold=True)
    for k, v in bundle["Recommendations"].items():
        add_line(f"• {k}: {v}")

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
